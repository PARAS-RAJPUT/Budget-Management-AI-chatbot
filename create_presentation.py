from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Format title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 112, 192)
    
    # Format subtitle
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(32)
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(89, 89, 89)

def create_content_slide(prs, title, content):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = title
    content_shape.text = content
    
    # Format title
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 112, 192)
    
    # Format content
    for paragraph in content_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = RGBColor(89, 89, 89)

def main():
    prs = Presentation()
    
    # Slide 1: Title
    create_title_slide(prs, 
        "Budget Management Chatbot",
        "AI-Powered Financial Assistant")
    
    # Slide 2: Project Objective
    objectives = """• Create an intelligent chatbot for personal finance management
• Provide personalized financial advice using AI
• Help users with:
  - Budget planning and tracking
  - Savings strategies
  - Investment recommendations
  - Debt management
  - Financial goal setting"""
    create_content_slide(prs, "Project Objective", objectives)
    
    # Slide 3: How It Works
    workflow = """1. User sends a message through web interface
2. BudgetBot processes using Google Gemini AI
3. System analyzes financial context
4. AI generates personalized advice
5. Response sent back to user
6. User profile updated automatically"""
    create_content_slide(prs, "How It Works", workflow)
    
    # Slide 4: Key Features
    features = """• Natural Language Processing
• Context-aware responses
• Persistent user profiles
• Real-time financial advice
• Budget tracking
• Investment recommendations
• Debt management strategies"""
    create_content_slide(prs, "Key Features", features)
    
    # Slide 5: Sample Interactions
    interactions = """Budget Analysis:
User: "I make 5000 per month"
Bot: "Here's your personalized budget:
• Essential expenses (50%): $2,500
• Discretionary spending (30%): $1,500
• Savings (20%): $1,000

Investment Advice:
User: "How should I invest my savings?"
Bot: "Based on your profile, consider:
• 60% in index funds
• 30% in bonds
• 10% in emergency fund""""
    create_content_slide(prs, "Sample Interactions", interactions)
    
    # Slide 6: Project Links
    links = """• GitHub Repository: https://github.com/yourusername/budget-chatbot
• Live Demo: https://budget-chatbot-demo.herokuapp.com
• Future Enhancements:
  - User authentication
  - Database integration
  - Mobile app version
  - Advanced analytics"""
    create_content_slide(prs, "Project Links & Future Plans", links)
    
    # Save the presentation
    prs.save('Budget_Management_Chatbot.pptx')

if __name__ == '__main__':
    main() 